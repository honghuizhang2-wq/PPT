#!/usr/bin/env python3
"""
东亚人类反焦虑平静指南 - 手帐风格 PPT
模拟手写笔记本：纸张底纹、便签贴纸、胶带、涂鸦、手写字体感
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

# ── 手帐配色 ──
PAPER_CREAM   = RGBColor(0xFD, 0xF8, 0xEF)   # 纸张底色
PAPER_LINED   = RGBColor(0xF0, 0xE6, 0xD6)   # 横线纸
PAPER_GRID    = RGBColor(0xF5, 0xF0, 0xE8)   # 方格纸
PAPER_KRAFT   = RGBColor(0xE8, 0xD5, 0xB7)   # 牛皮纸

WASHI_PINK    = RGBColor(0xFC, 0xD5, 0xCE)   # 粉色和纸胶带
WASHI_YELLOW  = RGBColor(0xFF, 0xF3, 0xB0)   # 黄色胶带
WASHI_GREEN   = RGBColor(0xC8, 0xE6, 0xC9)   # 绿色胶带
WASHI_BLUE    = RGBColor(0xBB, 0xDE, 0xFB)   # 蓝色胶带
WASHI_LAVENDER= RGBColor(0xE1, 0xBE, 0xE7)   # 薰衣草胶带

STICKY_YELLOW = RGBColor(0xFF, 0xF9, 0xC4)   # 便签黄
STICKY_PINK   = RGBColor(0xFC, 0xE4, 0xEC)   # 便签粉
STICKY_GREEN  = RGBColor(0xE8, 0xF5, 0xE9)   # 便签绿

INK_BLACK     = RGBColor(0x2C, 0x2C, 0x2C)   # 墨水黑
INK_BROWN     = RGBColor(0x5D, 0x40, 0x37)   # 墨水棕
INK_BLUE      = RGBColor(0x1A, 0x23, 0x7E)   # 钢笔蓝
INK_RED       = RGBColor(0xC6, 0x28, 0x28)   # 红笔
INK_GREEN     = RGBColor(0x2E, 0x7D, 0x32)   # 绿笔

HIGHLIGHT_YEL = RGBColor(0xFF, 0xF1, 0x76)   # 荧光黄
HIGHLIGHT_PNK = RGBColor(0xFF, 0x80, 0xAB)   # 荧光粉

prs = Presentation()
prs.slide_width  = Inches(9)
prs.slide_height = Inches(16)

# ── 工具函数 ──

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, w, h, fill_color, line_color=None, line_w=None):
    s = slide.shapes.add_shape(shape_type, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    return s

def add_rect(slide, left, top, w, h, fill, line=None, lw=None):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, w, h, fill, line, lw)

def add_rounded(slide, left, top, w, h, fill, line=None, lw=None):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill, line, lw)

def add_circle(slide, left, top, size, fill):
    return add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill)

def text_box(slide, left, top, w, h, text, size=16, color=INK_BLACK,
             bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei", spacing=1.5):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.line_spacing = Pt(size * spacing)
    return box

def sticky_note(slide, left, top, w, h, text, bg=STICKY_YELLOW, size=14, color=INK_BROWN):
    """便签贴纸效果"""
    s = add_rounded(slide, left, top, w, h, bg)
    # 便签阴影（用底层深色矩形模拟）
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(size * 1.6)
    return s

def washi_tape(slide, left, top, w, h, color, rotation=0):
    """和纸胶带效果（半透明矩形）"""
    s = add_rect(slide, left, top, w, h, color)
    # 设置半透明 - 通过 XML 直接操作
    try:
        spPr = s._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                a_elem = srgb.makeelement(qn('a:alpha'), {'val': '50000'})
                srgb.append(a_elem)
    except Exception:
        pass  # 透明度不是必须的
    s.line.fill.background()
    if rotation:
        s.rotation = rotation
    return s

def draw_lines(slide, left, top, w, count, spacing, color):
    """模拟笔记本横线"""
    for i in range(count):
        y = top + Emu(int(spacing * (i + 1)))
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, w, Pt(0.8))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

def draw_dots(slide, left, top, cols, rows, spacing, color):
    """模拟点阵纸"""
    for r in range(rows):
        for c in range(cols):
            x = left + Emu(int(spacing * c))
            y = top + Emu(int(spacing * r))
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Pt(3), Pt(3))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

def doodle_star(slide, cx, cy, size, color):
    """涂鸦星星（用文字代替）"""
    text_box(slide, cx - Inches(0.2), cy - Inches(0.2), Inches(0.5), Inches(0.5),
             "✦", size=size, color=color, align=PP_ALIGN.CENTER)

def doodle_heart(slide, cx, cy, size, color):
    text_box(slide, cx - Inches(0.15), cy - Inches(0.15), Inches(0.4), Inches(0.4),
             "♡", size=size, color=color, align=PP_ALIGN.CENTER)

def photo_corner(slide, x, y, size=Inches(0.3)):
    """照片角贴效果"""
    c = INK_BROWN
    # 左上角
    add_rect(slide, x, y, size, Pt(3), c)
    add_rect(slide, x, y, Pt(3), size, c)

def emoji_deco(slide, left, top, emoji, size=32):
    text_box(slide, left, top, Inches(0.6), Inches(0.6), emoji, size=size, align=PP_ALIGN.CENTER)

# ================================================================
# Slide 1: 手帐封面
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_CREAM)

# 顶部胶带装饰
washi_tape(slide, Inches(0.5), Inches(0.8), Inches(8.0), Inches(0.5), WASHI_PINK)
washi_tape(slide, Inches(1.0), Inches(1.5), Inches(7.0), Inches(0.4), WASHI_YELLOW, rotation=-2)

# 中心主标题区 - 模拟剪贴纸
add_rect(slide, Inches(1.0), Inches(3.5), Inches(7.0), Inches(6.0), PAPER_LINED, INK_BROWN, 1.5)
# 横线
draw_lines(slide, Inches(1.3), Inches(4.0), Inches(6.4), 12, Inches(0.45), RGBColor(0xD0, 0xC8, 0xB8))

# 标题手写感
text_box(slide, Inches(1.5), Inches(4.0), Inches(6.0), Inches(1.2),
         "东亚人类", size=52, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1.5), Inches(5.5), Inches(6.0), Inches(1.2),
         "反焦虑平静指南", size=44, color=INK_GREEN, bold=True, align=PP_ALIGN.CENTER)

# 副标题
text_box(slide, Inches(1.5), Inches(7.2), Inches(6.0), Inches(1.0),
         "如果你也经常觉得「还不够好」\n这本手帐是写给你的", size=18,
         color=INK_BROWN, align=PP_ALIGN.CENTER, spacing=1.8)

# 荧光笔高亮效果
add_rect(slide, Inches(2.5), Inches(8.5), Inches(4.0), Inches(0.4), HIGHLIGHT_YEL)
text_box(slide, Inches(2.5), Inches(8.4), Inches(4.0), Inches(0.5),
         "🌿 做一个平静的东亚人 🌿", size=14, color=INK_BROWN, align=PP_ALIGN.CENTER)

# 底部胶带
washi_tape(slide, Inches(0.3), Inches(10.0), Inches(8.4), Inches(0.4), WASHI_GREEN, rotation=1)

# 照片角贴
photo_corner(slide, Inches(1.0), Inches(3.5))
photo_corner(slide, Inches(7.3), Inches(3.5))
photo_corner(slide, Inches(1.0), Inches(9.2))
photo_corner(slide, Inches(7.3), Inches(9.2))

# 散落装饰
doodle_star(slide, Inches(0.5), Inches(2.5), 28, HIGHLIGHT_PNK)
doodle_heart(slide, Inches(8.0), Inches(3.0), 32, INK_RED)
emoji_deco(slide, Inches(7.5), Inches(2.0), "🌱", 36)
emoji_deco(slide, Inches(0.8), Inches(11.0), "☁️", 30)
doodle_star(slide, Inches(8.2), Inches(10.5), 24, HIGHLIGHT_YEL)

# ================================================================
# Slide 2: 焦虑清单 - 方格纸风格
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_GRID)

# 方格纸点阵
draw_dots(slide, Inches(0.3), Inches(0.3), 50, 90, Inches(0.18), RGBColor(0xCC, 0xCC, 0xCC))

# 顶部胶带 + 标题
washi_tape(slide, Inches(0.5), Inches(0.5), Inches(8.0), Inches(0.5), WASHI_BLUE)
text_box(slide, Inches(0.8), Inches(0.55), Inches(7.4), Inches(0.5),
         "📋 你的焦虑清单", size=32, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)

text_box(slide, Inches(1.0), Inches(1.3), Inches(7.0), Inches(0.5),
         "看看你中了几条？（全中的请举手🙋）", size=14, color=INK_BROWN, align=PP_ALIGN.CENTER)

# 焦虑项 - 便签贴纸风格
anxieties = [
    ("😰", "同龄人焦虑", "别人结婚生娃升职\n我还在想今天吃什么", STICKY_YELLOW),
    ("📱", "社交媒体焦虑", "朋友圈人均精致生活\n我连被子都没叠", STICKY_PINK),
    ("⏰", "年龄焦虑", "25岁觉得来不及了\n30岁觉得人生完了", STICKY_GREEN),
    ("💰", "金钱焦虑", "工资涨幅跑不过\n奶茶涨价速度", STICKY_YELLOW),
    ("📝", "自我审查", "发条朋友圈要P图1小时\n最后设为仅自己可见", STICKY_PINK),
    ("😔", "休息羞耻", "休息一天就焦虑\n觉得自己在浪费生命", STICKY_GREEN),
]

for i, (emoji, title, desc, bg) in enumerate(anxieties):
    row, col = divmod(i, 2)
    left = Inches(0.5 + col * 4.3)
    top = Inches(2.0 + row * 4.2)
    # 便签
    sticky_note(slide, left, top, Inches(3.8), Inches(3.8), "", bg, 14)
    # emoji
    emoji_deco(slide, left + Inches(1.4), top + Inches(0.2), emoji, 36)
    # 标题
    text_box(slide, left + Inches(0.3), top + Inches(0.9), Inches(3.2), Inches(0.5),
             title, size=18, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)
    # 描述
    text_box(slide, left + Inches(0.3), top + Inches(1.5), Inches(3.2), Inches(1.8),
             desc, size=13, color=INK_BROWN, align=PP_ALIGN.CENTER, spacing=1.7)
    # 便签上的胶带
    washi_tape(slide, left + Inches(0.8), top - Inches(0.1), Inches(2.2), Inches(0.3),
               [WASHI_PINK, WASHI_YELLOW, WASHI_GREEN, WASHI_BLUE][i % 4], rotation=(-3 + i * 2))

# 底部便签
sticky_note(slide, Inches(1.5), Inches(14.5), Inches(6.0), Inches(1.0),
            "💛 如果你中了3条以上 说明你是一个正常的东亚人", size=13)

# ================================================================
# Slide 3: 焦虑根源 - 牛皮纸剪贴簿
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_CREAM)

# 标题区 - 牛皮纸标签
add_rounded(slide, Inches(2.0), Inches(0.6), Inches(5.0), Inches(0.9), PAPER_KRAFT, INK_BROWN, 1)
text_box(slide, Inches(2.0), Inches(0.65), Inches(5.0), Inches(0.8),
         "🔍 焦虑从哪里来？", size=28, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)

roots = [
    ("🏫", "教育系统", [
        "从小被教育「不能输在起跑线」",
        "考了98分被问另外2分呢",
        "第一名只有一个",
        "但所有人都被要求争第一",
    ], WASHI_PINK),
    ("👨‍👩‍👧", "家庭期待", [
        "「别人家的孩子」是终身阴影",
        "听话是美德，叛逆是罪过",
        "你的选择永远不够好",
        "他们的担心永远不够多",
    ], WASHI_YELLOW),
    ("🌏", "社会环境", [
        "35岁危机、内卷、996",
        "「你不干有的是人干」",
        "成功学把人逼成机器",
        "却忘了问你快不快乐",
    ], WASHI_GREEN),
]

for i, (emoji, title, points, tape_color) in enumerate(roots):
    top = Inches(2.0 + i * 4.5)

    # 左侧胶带标签
    washi_tape(slide, Inches(0.5), top + Inches(0.3), Inches(0.5), Inches(3.2), tape_color)
    text_box(slide, Inches(0.3), top + Inches(1.0), Inches(0.9), Inches(1.0),
             emoji, size=36, align=PP_ALIGN.CENTER)

    # 右侧内容卡片 - 模拟剪贴纸
    card_bg = [STICKY_YELLOW, STICKY_PINK, STICKY_GREEN][i]
    add_rounded(slide, Inches(1.3), top, Inches(7.0), Inches(3.8), card_bg, INK_BROWN, 0.5)
    photo_corner(slide, Inches(1.3), top)
    photo_corner(slide, Inches(7.7), top)

    text_box(slide, Inches(1.6), top + Inches(0.2), Inches(6.4), Inches(0.5),
             title, size=22, color=INK_BLACK, bold=True)

    # 模拟手写列表
    for j, point in enumerate(points):
        y = top + Inches(0.9 + j * 0.65)
        text_box(slide, Inches(1.8), y, Inches(6.0), Inches(0.5),
                 f"• {point}", size=14, color=INK_BROWN, spacing=1.4)

# ================================================================
# Slide 4: 认知重启 - 横线笔记本
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_LINED)

# 横线
draw_lines(slide, Inches(0.5), Inches(0.3), Inches(8.0), 45, Inches(0.35), RGBColor(0xCC, 0xBB, 0xAA))

# 左侧红线（模拟笔记本装订线）
add_rect(slide, Inches(0.8), Inches(0), Pt(2), Inches(16), INK_RED)
add_rect(slide, Inches(0.9), Inches(0), Pt(0.5), Inches(16), INK_RED)

# 顶部便签
sticky_note(slide, Inches(1.2), Inches(0.5), Inches(7.0), Inches(1.0),
            "🧠 认知重启 · 换个脑子想问题", STICKY_YELLOW, 24, INK_BLACK)

tips = [
    ("💡", "允许自己「普通」", "世界上99%的人都是普通人\n普通不是失败\n是一种正常的人生状态", WASHI_PINK),
    ("🔄", "把「来不及」换成「刚刚好」", "摩西奶奶78岁才开始画画\n人生不是百米冲刺\n是各自赛道的马拉松", WASHI_YELLOW),
    ("🎯", "定义自己的成功", "成功不是只有一种模板\n今天开心了就是成功\n睡了个好觉也是成功", WASHI_GREEN),
    ("🚫", "屏蔽噪音", "别人的评价是他们的事\n你的人生是你的事\n学会「关你什么事」", WASHI_BLUE),
]

for i, (emoji, title, desc, tape) in enumerate(tips):
    row, col = divmod(i, 2)
    left = Inches(1.2 + col * 3.8)
    top = Inches(2.0 + row * 6.8)

    # 胶带固定效果
    washi_tape(slide, left + Inches(0.5), top - Inches(0.1), Inches(2.8), Inches(0.35), tape, rotation=(-2 + i))

    # 便签
    bg = [STICKY_PINK, STICKY_YELLOW, STICKY_GREEN, STICKY_PINK][i]
    sticky_note(slide, left, top, Inches(3.5), Inches(5.8), "", bg, 14)

    emoji_deco(slide, left + Inches(1.2), top + Inches(0.3), emoji, 40)
    text_box(slide, left + Inches(0.3), top + Inches(1.0), Inches(2.9), Inches(0.6),
             title, size=17, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, left + Inches(0.3), top + Inches(1.8), Inches(2.9), Inches(3.5),
             desc, size=13, color=INK_BROWN, align=PP_ALIGN.CENTER, spacing=1.8)

# ================================================================
# Slide 5: 行动指南 - 活页本风格
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_CREAM)

# 顶部胶带
washi_tape(slide, Inches(0.3), Inches(0.5), Inches(8.4), Inches(0.45), WASHI_LAVENDER)

text_box(slide, Inches(0.8), Inches(0.55), Inches(7.4), Inches(0.5),
         "🌿 行动指南 · 今天就能做的事", size=28, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)

actions = [
    ("1️⃣", "每天留10分钟给自己", "什么都不做，就发呆\n让大脑休息一下"),
    ("2️⃣", "取关让你焦虑的账号", "不舒服的内容就屏蔽\n你的信息流你做主"),
    ("3️⃣", "学会说「不知道」", "不用什么都知道\n「我不知道」是一种勇气"),
    ("4️⃣", "降低期待值", "对别人降低期待\n对自己也降低期待"),
    ("5️⃣", "找到你的「小确幸」", "一杯好喝的奶茶\n一个温暖的午后"),
    ("6️⃣", "允许自己休息", "休息不是偷懒\n你不是机器"),
]

for i, (num, title, desc) in enumerate(actions):
    row, col = divmod(i, 2)
    left = Inches(0.5 + col * 4.3)
    top = Inches(1.5 + row * 4.5)

    # 背景纸片（稍微旋转模拟随意贴）
    rotation = [-1, 1.5, -0.5, 2, -1.5, 0.8][i]
    bg = [STICKY_YELLOW, STICKY_PINK, STICKY_GREEN, STICKY_YELLOW, STICKY_PINK, STICKY_GREEN][i]
    s = add_rounded(slide, left, top, Inches(3.8), Inches(4.0), bg, INK_BROWN, 0.5)
    s.rotation = rotation

    # 编号圆形标签
    add_circle(slide, left + Inches(0.15), top + Inches(0.15), Inches(0.6),
               [WASHI_PINK, WASHI_GREEN, WASHI_BLUE, WASHI_YELLOW, WASHI_LAVENDER, WASHI_PINK][i])
    text_box(slide, left + Inches(0.15), top + Inches(0.2), Inches(0.6), Inches(0.5),
             num, size=20, align=PP_ALIGN.CENTER)

    text_box(slide, left + Inches(0.9), top + Inches(0.2), Inches(2.7), Inches(0.5),
             title, size=17, color=INK_BLACK, bold=True)
    text_box(slide, left + Inches(0.3), top + Inches(0.9), Inches(3.2), Inches(2.8),
             desc, size=13, color=INK_BROWN, spacing=1.7)

    # 底部胶带
    washi_tape(slide, left + Inches(0.5), top + Inches(3.6), Inches(2.8), Inches(0.3),
               [WASHI_GREEN, WASHI_BLUE, WASHI_PINK, WASHI_GREEN, WASHI_YELLOW, WASHI_BLUE][i])

# ================================================================
# Slide 6: 松弛感语录 - 拍立得风格
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_CREAM)

text_box(slide, Inches(0.8), Inches(0.5), Inches(7.4), Inches(0.8),
         "💬 东亚人的松弛感语录", size=30, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1.0), Inches(1.3), Inches(7.0), Inches(0.5),
         "建议截图保存，焦虑时拿出来读一读", size=13, color=INK_BROWN, align=PP_ALIGN.CENTER)

quotes = [
    ("「来得及，一切都来得及」", STICKY_YELLOW, WASHI_PINK),
    ("「不完美也没关系」", STICKY_PINK, WASHI_GREEN),
    ("「今天不努力，明天也可以」", STICKY_GREEN, WASHI_BLUE),
    ("「别人的成功跟我没关系」", STICKY_YELLOW, WASHI_LAVENDER),
    ("「我值得被爱，不需要条件」", STICKY_PINK, WASHI_YELLOW),
    ("「慢慢来，比较快」", STICKY_GREEN, WASHI_PINK),
    ("「人生不是竞赛，是体验」", STICKY_YELLOW, WASHI_BLUE),
    ("「我已经很棒了」", STICKY_PINK, WASHI_GREEN),
]

for i, (quote, bg, tape) in enumerate(quotes):
    row, col = divmod(i, 2)
    left = Inches(0.5 + col * 4.3)
    top = Inches(2.0 + row * 3.3)
    rotation = [-2, 1.5, -1, 2, -1.5, 1, -0.5, 2.5][i]

    # 拍立得白框
    frame = add_rect(slide, left, top, Inches(3.8), Inches(2.8), PAPER_CREAM, INK_BROWN, 1)
    frame.rotation = rotation

    # 内容区
    inner = add_rect(slide, left + Inches(0.15), top + Inches(0.15),
                     Inches(3.5), Inches(2.0), bg)
    inner.rotation = rotation

    # 胶带
    washi_tape(slide, left + Inches(0.8), top - Inches(0.15), Inches(2.2), Inches(0.3), tape)

    text_box(slide, left + Inches(0.3), top + Inches(0.5), Inches(3.2), Inches(1.5),
             quote, size=16, color=INK_BROWN, align=PP_ALIGN.CENTER, spacing=1.6)

# ================================================================
# Slide 7: 给你的信 - 手写信纸
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_CREAM)

# 信纸横线
draw_lines(slide, Inches(0.8), Inches(2.0), Inches(7.4), 30, Inches(0.4), RGBColor(0xCC, 0xBB, 0xAA))

# 左侧红线
add_rect(slide, Inches(0.7), Inches(0), Pt(2), Inches(16), INK_RED)

# 顶部胶带
washi_tape(slide, Inches(2.0), Inches(0.5), Inches(5.0), Inches(0.45), WASHI_PINK, rotation=-1)

text_box(slide, Inches(1.0), Inches(0.6), Inches(7.0), Inches(0.8),
         "🌱 最后想对你说", size=32, color=INK_BLACK, bold=True, align=PP_ALIGN.CENTER)

message_lines = [
    "亲爱的东亚人类，",
    "",
    "你不需要成为「别人家的孩子」",
    "你不需要在30岁之前完成所有事",
    "你不需要永远保持高效和完美",
    "",
    "你可以慢一点",
    "你可以普通一点",
    "你可以不那么「有用」",
    "",
    "你存在的意义",
    "不是为了满足所有人的期待",
    "而是为了体验这仅有一次的人生",
    "",
    "今天开始",
    "允许自己做一个",
    "平静而幸福的普通人 🌿",
]

for i, line in enumerate(message_lines):
    y = Inches(1.8 + i * 0.42)
    color = INK_BLACK if i == 0 else INK_BROWN
    bold = (i == 0)
    text_box(slide, Inches(1.2), y, Inches(6.5), Inches(0.4),
             line, size=16, color=color, bold=bold, align=PP_ALIGN.LEFT, spacing=1.4)

# 底部装饰
emoji_deco(slide, Inches(7.0), Inches(14.5), "🌸", 28)
doodle_heart(slide, Inches(1.5), Inches(14.8), 24, INK_RED)

# ================================================================
# Slide 8: 结尾 - 手帐封底
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, PAPER_KRAFT)

# 胶带装饰
washi_tape(slide, Inches(0.5), Inches(1.0), Inches(8.0), Inches(0.5), WASHI_PINK, rotation=1)
washi_tape(slide, Inches(1.0), Inches(1.5), Inches(7.0), Inches(0.4), WASHI_YELLOW, rotation=-1)

# 中心白纸
add_rect(slide, Inches(1.5), Inches(4.0), Inches(6.0), Inches(8.0), PAPER_CREAM, INK_BROWN, 1)
photo_corner(slide, Inches(1.5), Inches(4.0))
photo_corner(slide, Inches(6.9), Inches(4.0))
photo_corner(slide, Inches(1.5), Inches(11.4))
photo_corner(slide, Inches(6.9), Inches(11.4))

text_box(slide, Inches(2.0), Inches(5.0), Inches(5.0), Inches(1.0),
         "东亚人类", size=36, color=INK_GREEN, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(2.0), Inches(6.0), Inches(5.0), Inches(1.0),
         "反焦虑平静指南", size=36, color=INK_GREEN, bold=True, align=PP_ALIGN.CENTER)

# 荧光笔
add_rect(slide, Inches(2.5), Inches(7.5), Inches(4.0), Inches(0.4), HIGHLIGHT_YEL)
text_box(slide, Inches(2.5), Inches(7.4), Inches(4.0), Inches(0.5),
         "愿你平静，愿你自由，愿你幸福 🌿", size=15, color=INK_BROWN, align=PP_ALIGN.CENTER)

sticky_note(slide, Inches(2.5), Inches(8.5), Inches(4.0), Inches(1.5),
            "觉得有用就收藏吧 💛\n焦虑的时候拿出来看看", size=15)

# 散落装饰
doodle_star(slide, Inches(0.8), Inches(3.0), 28, HIGHLIGHT_PNK)
doodle_star(slide, Inches(8.0), Inches(3.5), 24, HIGHLIGHT_YEL)
emoji_deco(slide, Inches(0.3), Inches(13.5), "☁️", 30)
emoji_deco(slide, Inches(8.0), Inches(14.0), "🌱", 28)

# 底部胶带
washi_tape(slide, Inches(0.3), Inches(14.5), Inches(8.4), Inches(0.4), WASHI_GREEN, rotation=0.5)

# 保存
output_path = "/root/.openclaw/workspace/东亚人类反焦虑平静指南_手帐版.pptx"
prs.save(output_path)
print(f"✅ 手帐风 PPT 已保存: {output_path}")
