#!/usr/bin/env python3
"""
XHS Style PPT Generator
Creates Xiaohongshu-style presentations with warm colors, card layouts,
sticky notes, and emoji decorations.

Usage:
  python3 create_xhs_ppt.py --slides-json slides.json --output out.pptx
  python3 create_xhs_ppt.py --slides-json slides.json --palette cool --output out.pptx
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip3 install --break-system-packages python-pptx")
    sys.exit(1)


# ── Color Palettes ──

PALETTES = {
    "warm": {
        "bg_colors": ["#FDF5E6", "#D4EDDA", "#FFF3CD", "#FCE4EC", "#E8DAEF", "#FFE5D0", "#D1F2EB"],
        "sticky": "#FFF9C4",
        "text_dark": "#2D2D2D",
        "text_body": "#5D4E37",
        "accent": "#2E7D32",
        "card_bg": "#FFFFFF",
    },
    "cool": {
        "bg_colors": ["#E8F4FD", "#B3E5FC", "#C8E6C9", "#F3E5F5", "#E0F7FA", "#F1F8E9", "#FFF8E1"],
        "sticky": "#E3F2FD",
        "text_dark": "#1A237E",
        "text_body": "#37474F",
        "accent": "#1565C0",
        "card_bg": "#FFFFFF",
    },
    "pastel": {
        "bg_colors": ["#F8BBD0", "#FFCCBC", "#FFF9C4", "#D1C4E9", "#C8E6C9", "#B3E5FC", "#F3E5F5"],
        "sticky": "#FFF9C4",
        "text_dark": "#4A148C",
        "text_body": "#5D4037",
        "accent": "#AD1457",
        "card_bg": "#FFFFFF",
    },
    "earth": {
        "bg_colors": ["#F5E6CC", "#D5DBDB", "#E8D5B7", "#D5E8D4", "#EDE7D9", "#E0D5C1", "#D7CCC8"],
        "sticky": "#FFF8E1",
        "text_dark": "#3E2723",
        "text_body": "#5D4037",
        "accent": "#33691E",
        "card_bg": "#FFFFFF",
    },
}


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class XHSPptBuilder:
    def __init__(self, palette_name="warm"):
        p = PALETTES.get(palette_name, PALETTES["warm"])
        self.bg_colors = [hex_to_rgb(c) for c in p["bg_colors"]]
        self.sticky = hex_to_rgb(p["sticky"])
        self.text_dark = hex_to_rgb(p["text_dark"])
        self.text_body = hex_to_rgb(p["text_body"])
        self.accent = hex_to_rgb(p["accent"])
        self.card_bg = hex_to_rgb(p["card_bg"])

        self.prs = Presentation()
        self.prs.slide_width = Inches(9)
        self.prs.slide_height = Inches(16)

    def _bg(self, slide, idx):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_colors[idx % len(self.bg_colors)]

    def _rect(self, slide, left, top, w, h, color, rounded=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        s = slide.shapes.add_shape(shape_type, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _circle(self, slide, left, top, size, color):
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _text(self, slide, left, top, w, h, text, size=18, color=None, bold=False,
              align=PP_ALIGN.LEFT, spacing=1.5):
        if color is None:
            color = self.text_dark
        box = slide.shapes.add_textbox(left, top, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.alignment = align
        p.line_spacing = Pt(size * spacing)
        return box

    def _emoji(self, slide, left, top, emoji, size=36):
        self._text(slide, left, top, Inches(0.8), Inches(0.8), emoji,
                   size=size, align=PP_ALIGN.CENTER)

    def _decorations(self, slide):
        """Add random circle decorations to corners."""
        import random
        colors = self.bg_colors[:4]
        positions = [
            (0.3, 0.3, 1.0), (7.5, 0.5, 0.7), (0.5, 14.5, 0.8), (7.0, 14.0, 1.2),
        ]
        for (x, y, s), c in zip(positions, colors):
            self._circle(slide, Inches(x), Inches(y), Inches(s), c)

    def build_cover(self, data, idx):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, idx)
        self._decorations(slide)

        # Title block
        self._rect(slide, Inches(0.8), Inches(4.0), Inches(7.4), Inches(5.5), self.bg_colors[1])
        self._text(slide, Inches(1.2), Inches(4.5), Inches(6.6), Inches(1.5),
                   data.get("title", ""), size=44, bold=True, align=PP_ALIGN.CENTER)
        if data.get("subtitle"):
            self._text(slide, Inches(1.5), Inches(6.5), Inches(6.0), Inches(2.0),
                       data["subtitle"], size=20, color=self.text_body,
                       align=PP_ALIGN.CENTER, spacing=1.8)
        if data.get("tagline"):
            self._text(slide, Inches(2.0), Inches(10.5), Inches(5.0), Inches(0.8),
                       data["tagline"], size=16, color=self.text_body, align=PP_ALIGN.CENTER)

    def build_cards(self, data, idx):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, idx)

        self._text(slide, Inches(0.8), Inches(0.8), Inches(7.4), Inches(1.0),
                   data.get("section_title", ""), size=32, bold=True, align=PP_ALIGN.CENTER)

        cards = data.get("cards", [])
        cols = 2
        for i, card in enumerate(cards):
            row, col = divmod(i, cols)
            left = Inches(0.6 + col * 4.2)
            top = Inches(2.2 + row * 3.8)
            w, h = Inches(3.8), Inches(3.4)

            self._rect(slide, left, top, w, h, self.card_bg)
            self._text(slide, left + Inches(0.15), top + Inches(0.15),
                       Inches(0.6), Inches(0.6), card.get("emoji", "📌"), size=28,
                       align=PP_ALIGN.CENTER)
            self._text(slide, left + Inches(0.15), top + Inches(0.6),
                       w - Inches(0.3), Inches(0.5), card.get("title", ""),
                       size=16, bold=True)
            self._text(slide, left + Inches(0.15), top + Inches(1.0),
                       w - Inches(0.3), h - Inches(1.2), card.get("desc", ""),
                       size=12, color=self.text_body, spacing=1.6)

        # Bottom sticky note
        if data.get("footer_note"):
            self._build_sticky(slide, Inches(1.5), Inches(14.2), Inches(6.0), Inches(1.2),
                               data["footer_note"], 14)

    def build_tips(self, data, idx):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, idx)

        self._text(slide, Inches(0.8), Inches(0.8), Inches(7.4), Inches(1.0),
                   data.get("section_title", ""), size=32, bold=True, align=PP_ALIGN.CENTER)

        tips = data.get("tips", [])
        cols = 2
        for i, tip in enumerate(tips):
            row, col = divmod(i, cols)
            left = Inches(0.6 + col * 4.2)
            top = Inches(2.0 + row * 4.5)
            w, h = Inches(3.8), Inches(4.0)

            self._rect(slide, left, top, w, h, self.card_bg)
            self._text(slide, left + Inches(0.15), top + Inches(0.15),
                       Inches(0.6), Inches(0.5), tip.get("num", "1️⃣"), size=22,
                       align=PP_ALIGN.CENTER)
            self._text(slide, left + Inches(0.15), top + Inches(0.6),
                       w - Inches(0.3), Inches(0.5), tip.get("title", ""),
                       size=17, bold=True)
            self._text(slide, left + Inches(0.15), top + Inches(1.2),
                       w - Inches(0.3), h - Inches(1.4), tip.get("desc", ""),
                       size=13, color=self.text_body, spacing=1.7)

    def build_quotes(self, data, idx):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, idx)

        self._text(slide, Inches(0.8), Inches(0.8), Inches(7.4), Inches(1.0),
                   data.get("section_title", ""), size=32, bold=True, align=PP_ALIGN.CENTER)

        if data.get("hint"):
            self._text(slide, Inches(1.0), Inches(1.8), Inches(7.0), Inches(0.6),
                       data["hint"], size=14, color=self.text_body, align=PP_ALIGN.CENTER)

        quotes = data.get("quotes", [])
        colors = [self.bg_colors[1], self.bg_colors[3], self.bg_colors[0], self.bg_colors[4]]
        for i, q in enumerate(quotes):
            row, col = divmod(i, 2)
            left = Inches(0.6 + col * 4.2)
            top = Inches(2.6 + row * 3.2)
            bg = colors[i % len(colors)]
            self._build_sticky(slide, left, top, Inches(3.8), Inches(2.6), q, 17, bg)

    def _build_sticky(self, slide, left, top, w, h, text, size=14, bg=None):
        if bg is None:
            bg = self.sticky
        s = self._rect(slide, left, top, w, h, bg)
        tf = s.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = self.text_body
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = Pt(size * 1.6)

    def build_message(self, data, idx):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, idx)
        self._decorations(slide)

        self._rect(slide, Inches(1.0), Inches(3.0), Inches(7.0), Inches(10.0), self.card_bg)
        self._text(slide, Inches(1.5), Inches(3.5), Inches(6.0), Inches(1.0),
                   data.get("title", ""), size=36, bold=True, align=PP_ALIGN.CENTER)
        self._text(slide, Inches(1.5), Inches(4.8), Inches(6.0), Inches(7.5),
                   data.get("body", ""), size=18, color=self.text_body,
                   align=PP_ALIGN.CENTER, spacing=1.7)

    def build_closing(self, data, idx):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, idx)
        self._decorations(slide)

        self._text(slide, Inches(1.0), Inches(5.0), Inches(7.0), Inches(1.5),
                   data.get("title", ""), size=36, bold=True, color=self.accent,
                   align=PP_ALIGN.CENTER)
        if data.get("tagline"):
            self._text(slide, Inches(1.5), Inches(7.0), Inches(6.0), Inches(1.0),
                       data["tagline"], size=20, color=self.text_body, align=PP_ALIGN.CENTER)
        if data.get("cta"):
            self._build_sticky(slide, Inches(2.0), Inches(9.0), Inches(5.0), Inches(1.5),
                               data["cta"], 16)

    def build_custom_section(self, data, idx):
        """Generic section with left-icon right-content layout."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, idx)

        self._text(slide, Inches(0.8), Inches(0.8), Inches(7.4), Inches(1.0),
                   data.get("section_title", ""), size=32, bold=True, align=PP_ALIGN.CENTER)

        items = data.get("items", [])
        for i, item in enumerate(items):
            top = Inches(2.2 + i * 4.5)
            # Left icon block
            self._rect(slide, Inches(0.8), top, Inches(1.5), Inches(3.8),
                       self.bg_colors[i % len(self.bg_colors)])
            self._text(slide, Inches(0.8), top + Inches(0.8), Inches(1.5), Inches(1.0),
                       item.get("emoji", ""), size=44, align=PP_ALIGN.CENTER)
            # Right content block
            self._rect(slide, Inches(2.5), top, Inches(5.7), Inches(3.8), self.card_bg)
            self._text(slide, Inches(2.8), top + Inches(0.2), Inches(5.0), Inches(0.6),
                       item.get("title", ""), size=22, bold=True)
            self._text(slide, Inches(2.8), top + Inches(0.9), Inches(5.0), Inches(2.6),
                       item.get("desc", ""), size=14, color=self.text_body, spacing=1.7)

    def build(self, slides_data, output_path):
        builders = {
            "cover": self.build_cover,
            "cards": self.build_cards,
            "tips": self.build_tips,
            "quotes": self.build_quotes,
            "message": self.build_message,
            "closing": self.build_closing,
            "section": self.build_custom_section,
        }

        for i, slide_data in enumerate(slides_data):
            slide_type = slide_data.get("type", "cards")
            builder = builders.get(slide_type, self.build_cards)
            builder(slide_data, i)

        self.prs.save(output_path)
        print(f"✅ PPT saved: {output_path} ({len(slides_data)} slides)")


def main():
    parser = argparse.ArgumentParser(description="XHS Style PPT Generator")
    parser.add_argument("--slides-json", required=True, help="Path to slides JSON file")
    parser.add_argument("--output", default="xhs_output.pptx", help="Output PPTX path")
    parser.add_argument("--palette", default="warm", choices=PALETTES.keys(),
                        help="Color palette (default: warm)")
    args = parser.parse_args()

    with open(args.slides_json, "r", encoding="utf-8") as f:
        slides_data = json.load(f)

    builder = XHSPptBuilder(palette_name=args.palette)
    builder.build(slides_data, args.output)


if __name__ == "__main__":
    main()
